import sys
import os
import shutil

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

# Target output directory
target_dir = r'd:\26_강의자료\프로젝트\발표용 공유'
os.makedirs(target_dir, exist_ok=True)
print(f"Target directory verified: {target_dir}")

# Build PPTX File using python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BLUE = RGBColor(15, 23, 42)     # #0F172A
ROYAL_BLUE = RGBColor(37, 99, 235)   # #2563EB
SKY_BLUE = RGBColor(56, 189, 248)    # #38BDF8
SLATE_GRAY = RGBColor(148, 163, 184) # #94A3B8
WHITE = RGBColor(255, 255, 255)
CARD_BG = RGBColor(30, 41, 59)       # #1E293B
GREEN = RGBColor(52, 211, 153)       # #34D399
AMBER = RGBColor(251, 191, 36)      # #FBBF24

def set_slide_background(slide, color=DARK_BLUE):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, tag_text, title_text, slide_num_str):
    # Tag
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = tag_text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SKY_BLUE
    
    # Title
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title_text
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    
    # Slide Num
    txBox3 = slide.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.2), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = slide_num_str
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = SLATE_GRAY
    p3.alignment = PP_ALIGN.RIGHT

# SLIDE 1: Title Slide
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1, DARK_BLUE)

shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11), Inches(1.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "FINAL PRESENTATION · 3040 HOUSING MODEL"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = SKY_BLUE

txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "3040 맞벌이를 위한 서울 최적 아파트 주거지 분석"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(11), Inches(1.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = '"내 예산으로 출퇴근도 편하고, 하락장에도 집값이 덜 떨어지는 서울 동네는 어디일까?"'
p.font.size = Pt(18)
p.font.color.rgb = SLATE_GRAY

txBox = slide1.shapes.add_textbox(Inches(1.2), Inches(5.5), Inches(11), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "• v2 실거래 정제본 83,178건 완벽 연동   • 머신러닝 회귀 설명력 R² 83.8%   • MCDM TOPSIS & K-Means 4대 세그먼트"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = GREEN

# Helper for card container
def create_card(slide, left, top, width, height, title, content_list, header_color=SKY_BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(51, 65, 85)
    
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = header_color
    
    for item in content_list:
        p2 = tf.add_paragraph()
        p2.text = f"• {item}"
        p2.font.size = Pt(12.5)
        p2.font.color.rgb = SLATE_GRAY
        p2.space_before = Pt(6)

# SLIDE 2: Problem Definition
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "01. PROBLEM DEFINITION", "3040 맞벌이 주거지의 3대 핵심 고민 및 목표", "02 / 10")

create_card(slide2, 0.8, 1.8, 3.6, 4.8, "01. 예산의 한계", [
    "6억 ~ 15억 원대 가용 예산",
    "실제 매수 가능한 아파트 단지 필터링",
    "수치 왜곡 방지를 위한 중앙 매매가 직접 사용"
], SKY_BLUE)

create_card(slide2, 4.8, 1.8, 3.6, 4.8, "02. 맞벌이 출퇴근 부담", [
    "부부의 직장 위치가 서로 다를 때 절충 필요",
    "GBD, YBD, CBD 3대 업무지구 20분대 도달",
    "07~09시 피크 타임 실측 이동량 적용"
], GREEN)

create_card(slide2, 8.8, 1.8, 3.6, 4.8, "03. 자산 하락 위험", [
    "금리 인상기 하락장 낙폭율 측정",
    "하락기(-7~-12%) 가격 방어력 최상 동네 선별",
    "금리 인하기 고탄력 반등 지역 포착"
], AMBER)

# SLIDE 3: Data Integration & v2 Cleanup
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "02. DATA INTEGRATION & CLEANUP", "v2 실거래 정제본(83,178건) 연동 및 노이즈 전처리", "03 / 10")

create_card(slide3, 0.8, 1.8, 5.6, 4.8, "✅ 필수 결합 데이터셋", [
    "실거래가 정제본 v2 (11_실거래가_정제본_v2.csv 83,178건)",
    "아파트 유효 매매 46,644건 정밀 추출",
    "동네별 아파트 평균 건축년도 (연식) 연동",
    "실거래 기반 주요 대표 아파트 단지명 직결 (래미안, 자이 등)",
    "07~09시 3040 실측 통근 이동량 O-D 결합"
], GREEN)

create_card(slide3, 6.8, 1.8, 5.6, 4.8, "❌ 제거된 노이즈 데이터", [
    "상권 분석 서비스 (점포수, 길단위 인구) 전량 삭제",
    "빌라, 단독/다세대, 오피스텔 비주거 주택 제외",
    "30억 초과 초고가 단지 및 100세대 미만 나홀로 아파트 제외",
    "심야 및 주말 비통근 이동량 제외"
], RGBColor(248, 113, 113))

# SLIDE 4: Regression Modeling
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "03. REGRESSION MODELING", "회귀 모델 분석: 무엇이 3040 주거 만족도를 결정하는가?", "04 / 10")

# Stat Card
card_stat = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(4.5), Inches(4.8))
card_stat.fill.solid()
card_stat.fill.fore_color.rgb = CARD_BG
card_stat.line.color.rgb = ROYAL_BLUE
tf_stat = card_stat.text_frame
p = tf_stat.paragraphs[0]
p.text = "R² = 83.8%"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = SKY_BLUE
p.alignment = PP_ALIGN.CENTER

p2 = tf_stat.add_paragraph()
p2.text = "Random Forest 회귀 예측 설명력\n(RMSE 오차: 3.55점)"
p2.font.size = Pt(14)
p2.font.color.rgb = SLATE_GRAY
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(10)

create_card(slide4, 5.7, 1.8, 6.8, 4.8, "📊 특성 중요도 (Feature Importance Top 4)", [
    "1위: 3대 도심 평균 통근시간 (64.5% - 압도적 1위)",
    "2위: YBD(여의도) 직주근접 소요시간 (12.7%)",
    "3위: 3040세대 상주인구 비중 (8.6%)",
    "4위: 30대 인구 비율 (3.0%)",
    "💡 인사이트: 주거 만족도의 77% 이상은 '직주근접 및 통근시간 단축'에 기인함"
], SKY_BLUE)

# SLIDE 5: TOPSIS Recommendation
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "04. RECOMMENDATION 1: TOPSIS MODEL", "개인 맞춤형 TOPSIS 의사결정 알고리즘", "05 / 10")

create_card(slide5, 0.8, 1.8, 3.6, 4.8, "A. 신혼 & 첫 집 (가성비)", [
    "가중치: 통근50% / 예산30%",
    "1위: 영등포구 영등포동2가 (7.35억)",
    "2위: 종로구 숭인동 (5.20억)",
    "3위: 관악구 남현동 (7.22억)",
    "특징: 5~7억대 가성비 실속형"
], SKY_BLUE)

create_card(slide5, 4.8, 1.8, 3.6, 4.8, "B. 자산방어형 (하락장 방어)", [
    "가중치: 방어50% / 통근30%",
    "1위: 영등포구 영등포동2가 (+1.5%)",
    "2위: 종로구 숭인동 (+1.3%)",
    "3위: 구로구 구로동 (-2.4%)",
    "특징: 하락장 가격 방어력 최상"
], GREEN)

create_card(slide5, 8.8, 1.8, 3.6, 4.8, "C. 직주근접 & 도심성장형", [
    "가중치: 통근40% / 3040 30%",
    "1위: 강서구 염창동 (43.3%)",
    "2위: 강서구 가양동 (43.3%)",
    "3위: 영등포구 당산동 (38.5%)",
    "특징: 3040 밀집도 및 출퇴근 1위"
], AMBER)

# SLIDE 6: K-Means 4 Clusters
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "05. RECOMMENDATION 2: K-MEANS CLUSTERING", "K-Means 4대 주거지 세그먼트 군집 분석", "06 / 10")

create_card(slide6, 0.8, 1.8, 5.6, 2.2, "그룹 1: 가성비 실속 & 실수요 밀집지", [
    "시세: 6.0억 ~ 8.3억 원 | 통근: 20~25분",
    "대표동네: 강서 염창·등촌, 관악 봉천 (3040비중 43.3% 서울 1위)"
], SKY_BLUE)

create_card(slide6, 6.8, 1.8, 5.6, 2.2, "그룹 2: 광역 허브 & 철통 방어선 🏆", [
    "시세: 9.6억 ~ 11.5억 원 | 하락장 낙폭: -7.1% ~ -12.0% (방어 1위)",
    "대표동네: 영등포 당산, 동작 사당, 구로 신도림 (실수요 추천 1위)"
], GREEN)

create_card(slide6, 0.8, 4.4, 5.6, 2.2, "그룹 3: 쿼드러플 & 도심 거점 🚀", [
    "시세: 12.1억 ~ 15.4억 원 | 상승기 반등률: +40.9% (고탄력)",
    "대표동네: 마포 공덕·도화, 송파 문정·가락 (4개 노선 쿼드러플)"
], AMBER)

create_card(slide6, 6.8, 4.4, 5.6, 2.2, "그룹 4: 자녀 보육 & 학군 배후지 🎓", [
    "시세: 7.5억 ~ 10.5억 원 | 통근: 30~40분 (주간인구 70%대)",
    "대표동네: 노원 중계(은행사거리), 강동 고덕·명일 (서울 3대 학군)"
], RGBColor(192, 132, 252))

# SLIDE 7: TOP 4 Matrix
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_header(slide7, "06. TOP NEIGHBORHOOD MATRIX", "서울 대표 추천 법정동 TOP 4 실측 매트릭스 (v2 연동)", "07 / 10")

rows, cols = 5, 6
left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8)
table_shape = slide7.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

headers = ["법정동", "중앙 매매가", "v2 연동 대표 단지", "통근시간", "하락장 낙폭", "추천 포인트"]
for col_idx, text in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = ROYAL_BLUE
    p = cell.text_frame.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE

matrix_data = [
    ["영등포구 당산동", "11.5억 원", "당산 삼성래미안, 삼환", "YBD 5분/GBD 15분", "-12.0% (방어1위)", "맞벌이 1순위 추천, 2·9호선 환승 요충지"],
    ["동작구 사당동", "10.7억 원", "사당 래미안로이뷰, 우성3차", "GBD 10분/CBD 20분", "-7.1% (서울최저)", "자산방어 1위, 강남/도심 20분 사통팔달"],
    ["강서구 염창동", "8.3억 원", "염창 래미안, 동아3차", "YBD 10분/GBD 25분", "-29.1%", "6~8억 가성비 1등, 9호선 급행 역세권"],
    ["마포구 공덕동", "14.1억 원", "공덕 래미안4차, 공덕자이", "YBD 5분/CBD 10분", "+40.9% (상승기)", "직주근접 종결지, 4개 노선 쿼드러플 환승"]
]

for row_idx, row_data in enumerate(matrix_data):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11.5)
        p.font.color.rgb = WHITE if col_idx == 0 else SLATE_GRAY
        if col_idx == 4:
            p.font.color.rgb = GREEN

# SLIDE 8: Data QA & Verification
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8)
add_header(slide8, "07. DATA QA & VERIFICATION", "데이터 품질 검증 및 전처리 체크리스트", "08 / 10")

create_card(slide8, 0.8, 1.8, 5.6, 4.8, "✔️ 수치 검증 및 전처리 완료 사항", [
    "v2 실거래 정제본 83,178건 기반 아파트 46,644건 결합",
    "단가×면적 착시 제거를 위한 최근 중앙 매매가 직접 사용",
    "07:00~09:59 아침 피크 타임 실측 통근량 정밀 필터링",
    "다공선성 해결을 위한 Ridge Regression 정규화 적용",
    "K-Means 난수 seed=42 고정으로 결과 재현성 보장"
], GREEN)

create_card(slide8, 6.8, 1.8, 5.6, 4.8, "⚠️ 해석 시 유의사항 (Limitations)", [
    "과거 하락장 방어율이 미래 시세를 100% 보장하지는 않음",
    "순수 아파트 전용 분석이므로 빌라/오피스텔 시장과 구분 필요",
    "개별 아파트 단지의 층/향/연식은 법정동 평균으로 1차 스크리닝"
], AMBER)

# SLIDE 9: Dashboard & Roadmap
slide9 = prs.slides.add_slide(blank_layout)
set_slide_background(slide9)
add_header(slide9, "08. DASHBOARD & ROADMAP", "인터랙티브 웹 대시보드 및 향후 공유 로드맵", "09 / 10")

create_card(slide9, 0.8, 1.8, 5.6, 4.8, "🌐 인터랙티브 지도 대시보드", [
    "3040_맞벌이_주거지_지도_대시보드.html 탑재",
    "Leaflet GIS 기반 동네 마커 핀 시각화",
    "소형(59㎡) / 국평(84㎡) / 중대형 평형대 선택 필터",
    "동네 클릭 시 v2 실거래 아파트 단지명 및 연식 즉시 출력"
], SKY_BLUE)

create_card(slide9, 6.8, 1.8, 5.6, 4.8, "🚀 공유 폴더 구성 완료", [
    "D:\\26_강의자료\\프로젝트\\발표용 공유 폴더 내 일괄 정리",
    "한눈에 보는 요약 (HTML/MD)",
    "초간단 3분 동네 추천 가이드 (HTML/MD)",
    "발표용 PPTX 슬라이드 덱 및 1:1 맞춤 대본 포함"
], GREEN)

# SLIDE 10: Conclusion & Q&A
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)
add_header(slide10, "09. CONCLUSION & Q&A", "프로젝트 결론 및 질의응답", "10 / 10")

txBox = slide10.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(11), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "3040 맞벌이를 위한 최적 주거지 선택의 핵심은\n[가용 예산 범위] 내에서 [통근시간 단축]과 [하락장 자산방어율]의 균형을 맞추는 것입니다."
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox = slide10.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "경청해 주셔서 감사합니다 (Q & A)"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = SKY_BLUE
p.alignment = PP_ALIGN.CENTER

# Save PPTX file
pptx_path = os.path.join(target_dir, "3040_맞벌이_주거지_분석_발표자료.pptx")
prs.save(pptx_path)
print(f"PPTX file created at: {pptx_path}")

# Copy all final shareable files to '발표용 공유' folder
source_dir = r'd:\26_강의자료\프로젝트'
files_to_copy = [
    "3040_맞벌이_주거지_지도_대시보드.html",
    "3040_맞벌이_주거지_분석_한눈에보는_요약.html",
    "3040_맞벌이_쉬운_동네추천_가이드.html",
    "3040_맞벌이_주거지_분석_발표자료.html",
    "3040_맞벌이_주거지_분석_및_모델링_팀공유리포트.html",
    "3040_맞벌이_주거지_분석_발표대본_및_슬라이드.md",
    "3040_맞벌이_주거지_분석_한눈에보는_요약.md",
    "3040_맞벌이_쉬운_동네추천_가이드.md",
    "3040_맞벌이_주거지_클러스터링_분석결과_v2연동.csv"
]

copied_count = 0
for fname in files_to_copy:
    src = os.path.join(source_dir, fname)
    dst = os.path.join(target_dir, fname)
    if os.path.exists(src):
        shutil.copy2(src, dst)
        copied_count += 1
        print(f"Copied: {fname} -> {target_dir}")

print(f"\nSuccessfully organized {copied_count} files and 1 PPTX presentation deck in '{target_dir}'!")
