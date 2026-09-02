import sys
import os

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

TARGET_DIR = r'd:\26_강의자료\프로젝트\발표용 공유\최종'
os.makedirs(TARGET_DIR, exist_ok=True)
pptx_path = os.path.join(TARGET_DIR, '3040_맞벌이_주거지_분석_최종발표.pptx')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(15, 23, 42)
    BLUE_ACCENT = RGBColor(59, 130, 246)
    SKY_BLUE = RGBColor(56, 189, 248)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)
    CARD_BG = RGBColor(30, 41, 59)
    GREEN = RGBColor(52, 211, 153)
    GOLD = RGBColor(251, 191, 36)

    def set_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_header(slide, tag, title, slide_num):
        tb1 = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        p1 = tb1.text_frame.paragraphs[0]
        p1.text = tag
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = SKY_BLUE

        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10), Inches(0.8))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

        tb3 = slide.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.2), Inches(0.4))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = f"{slide_num} / 7"
        p3.font.size = Pt(12)
        p3.font.color.rgb = GRAY
        p3.alignment = PP_ALIGN.RIGHT

    def add_table(slide, left, top, width, height, headers, data):
        rows = len(data) + 1
        cols = len(headers)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            p = cell.text_frame.paragraphs[0]
            p.text = header
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = SKY_BLUE

        for row_idx, row_data in enumerate(data):
            for col_idx, val in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
                p = cell.text_frame.paragraphs[0]
                p.text = str(val)
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE

    blank_layout = prs.slide_layouts[6]

    # Slide 1: Title
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.0))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "3040 맞벌이를 위한 서울 아파트 주거지 분석"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p_sub = tf.add_paragraph()
    p_sub.text = "통합 리포트 1:1 연동 발표 슬라이드 (v3.5)"
    p_sub.font.size = Pt(19)
    p_sub.font.color.rgb = SKY_BLUE

    # Slide 2: PEA
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "Ⅰ. 문제 정의", "우리 팀의 프로젝트 핵심 문제 정의 (PEA)", "2")
    headers_s2 = ["구분", "핵심 정의 내용", "기대 효과"]
    data_s2 = [
        ["📌 Problem", "시세/호재 위주 부동산 정보로 내 예산 맞춤 통근/안정성 탐색 난항", "통합 입지지표 도입"],
        ["🎯 Target", "가용 예산 6억~15억 원대 3040 실수요 맞벌이 부부", "실수요 1:1 커스텀 타깃"],
        ["🚀 Solution", "통근(35%) + 가격방어율(30%) + 3040비중(25%) 다기준 점수화", "의사결정시간 대폭 단축"]
    ]
    add_table(s2, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), headers_s2, data_s2)

    # Slide 3: 산출 공식
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "Ⅱ. 점수 산출 공식", "종합점수(TOPSIS) 산출 공식 및 당산동 58.8점 스펙", "3")
    headers_s3 = ["평가 지표", "가중치 비중", "당산동 점수", "점수 산정 로직"]
    data_s3 = [
        ["🚀 통근 편리성", "35%", "88.4점", "직장 30분 이내 100점 만점! (YBD 5분 / GBD 15분)"],
        ["🛡️ 가격 방어력", "30%", "89.7점", "하락장 낙폭 -10% 이내 100점 만점! (-12.0% 철통 방어)"],
        ["👨‍👩‍👧 3040 친화도", "25%", "60.8점", "3040 인구 비중 35% 이상 100점 만점! (3040 비중 38.5%)"],
        ["💧 거래 유동성", "10%", "34.0점", "연간 아파트 매매 거래량 환금성 산출"],
        ["🏆 종합점수 (TOPSIS)", "100%", "58.8점", "서울시 137개 분석 법정동 중 종합 1위 산출!"]
    ]
    add_table(s3, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), headers_s3, data_s3)

    # Slide 4: 4대 지표 문턱값
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "Ⅲ. 문턱값 산정", "4대 지표별 문턱값(Threshold) 및 서울시 43.8분 앵커", "4")
    headers_s4 = ["평가 축", "문턱값(Threshold) 100점 만점 기준", "기준선 (Anchor)"]
    data_s4 = [
        ["1. 통근 편리성", "출퇴근 30분 이내는 전부 100점 만점 부여", "서울시 평균 출퇴근시간 (43.8분)"],
        ["2. 가격 방어력", "금리 폭락장 낙폭 -10% 이내는 100점 만점", "서울시 하락장 최저 낙폭선"],
        ["3. 3040 친화도", "3040대 실수요 인구 비중 35% 이상 100점 만점", "서울시 평균 상주 인구비"],
        ["4. 거래 유동성", "연간 매매 거래건수 50건 이상 100점 만점", "최소 환금성 보장선"]
    ]
    add_table(s4, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), headers_s4, data_s4)

    # Slide 5: 시나리오
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "Ⅳ. 맞춤 랭킹", "맞벌이 3대 커플 직장 조합 시나리오 추천", "5")
    headers_s5 = ["직장 조합", "추천 1위 동네", "중앙가", "부부 평균 통근시간", "핵심 추천 사유"]
    data_s5 = [
        ["강남 × 광화문", "서대문구 연희동", "7.2억 원", "GBD 35분 / CBD 15분", "양 직장 중간 지점 가성비 최상위 1위"],
        ["강남 × 여의도", "영등포구 당산동", "11.0억 원", "GBD 15분 / YBD 5분", "2·9호선 급행 환승 초역세권 라인"],
        ["여의도 × 광화문", "마포구 공덕동", "14.1억 원", "YBD 5분 / CBD 10분", "5호선 직결 쿼드러플 환승 교통 중심지"]
    ]
    add_table(s5, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), headers_s5, data_s5)

    # Slide 6: K-Means 군집 태그
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "Ⅴ. 주거 세그먼트", "K-Means 4대 주거지 세그먼트 (동네 유형 태그)", "6")
    headers_s6 = ["군집 세그먼트", "평균 시세", "통근시간", "하락장 방어력", "대표 동네"]
    data_s6 = [
        ["그룹 1: 가성비 실속형", "6.0억 ~ 8.3억", "20~25분", "-21.7% ~ -29.1%", "강서 염창 / 관악 봉천"],
        ["그룹 2: 광역 허브 ★", "9.6억 ~ 11.5억", "15~20분", "-7.1% ~ -12.0%", "영등포 당산 / 동작 사당 (추천 1위)"],
        ["그룹 3: 쿼드러플 도심", "12.1억 ~ 15.4억", "5~10분", "상승기 +40.9%", "마포 공덕 / 송파 문정"],
        ["그룹 4: 학군 배후지", "7.5억 ~ 10.5억", "30~40분", "-22.8% ~ -29.5%", "노원 중계 / 강동 고덕"]
    ]
    add_table(s6, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), headers_s6, data_s6)

    # Slide 7: 결론
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "Ⅵ. 최종 결론", "3040 맞벌이를 위한 서울 최적 동네 1위 당산동", "7")
    headers_s7 = ["평가 항목", "세부 수치", "서울시 대비 우수성"]
    data_s7 = [
        ["통근 소요시간", "YBD 5분 / GBD 15분", "2·9호선 환승 초역세권 (서울 최고 수준)"],
        ["금리충격 낙폭", "-12.0%", "서울 전체 최상위 자산 방어력 입증"],
        ["3040 거주 비중", "38.5%", "실수요 및 육아 커뮤니티 우수"],
        ["종합 랭킹 점수", "58.8점 (1위)", "실시간 가중치 조절 대시보드 튜닝 보장"]
    ]
    add_table(s7, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), headers_s7, data_s7)

    prs.save(pptx_path)
    print(f"PPTX saved: {pptx_path}")
except Exception as e:
    print(f"PPTX error: {e}")
